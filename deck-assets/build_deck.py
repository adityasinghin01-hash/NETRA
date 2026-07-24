#!/usr/bin/env python3
"""Build NETRA_deck_draft.pptx — a content-filled draft deck.

NOT the submission file. The KSP/Hack2Skill official template is what gets
submitted; this exists so the presentation teammate copies finished slides
across instead of retyping from markdown. Content mirrors docs/PPT_HANDBOOK.md.

Run:  python3 deck-assets/build_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "NETRA_deck_draft.pptx"

# NETRA dark theme
BG = RGBColor(0x0B, 0x12, 0x20)
INK = RGBColor(0xE2, 0xE8, 0xF0)
DIM = RGBColor(0x94, 0xA3, 0xB8)
MUTE = RGBColor(0x64, 0x74, 0x8B)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
OK = RGBColor(0x4A, 0xDE, 0x80)

W, H = Inches(13.333), Inches(7.5)  # 16:9


def new_deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slide(prs):
    """Blank slide with the dark background painted in."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, spacing=1.15):
    """runs = [(string, size_pt, bold, color), ...]; each is its own paragraph."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, size, bold, color) in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        r = para.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Inter"
    return tb


def heading(s, title, sub=None):
    runs = [(title, 34, True, INK)]
    if sub:
        runs.append((sub, 15, False, DIM))
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(1.3), runs)


def accent_bar(s, y=Inches(1.62), color=CYAN):
    bar = s.shapes.add_shape(1, Inches(0.7), y, Inches(1.5), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def picture(s, name, x, y, w):
    """Place a screenshot if present; otherwise leave a labelled placeholder."""
    p = HERE / name
    if p.exists():
        return s.shapes.add_picture(str(p), x, y, width=w)
    ph = s.shapes.add_shape(1, x, y, w, Inches(3))
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0x11, 0x1A, 0x2E)
    ph.line.color.rgb = MUTE
    tf = ph.text_frame
    tf.text = f"[ {name} missing ]"
    tf.paragraphs[0].runs[0].font.size = Pt(12)
    tf.paragraphs[0].runs[0].font.color.rgb = MUTE
    return ph


def note(s, body):
    """Speaker notes — the guidance that must not end up on the slide."""
    s.notes_slide.notes_text_frame.text = body


def footer(s, txt="Prototype · synthetic data generated to KSP's official FIR schema"):
    text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
         [(txt, 9, False, MUTE)])


def bullets(s, x, y, w, items, size=13, gap=0.52):
    for i, (head, body) in enumerate(items):
        yy = y + Inches(i * gap)
        text(s, x, yy, w, Inches(0.5),
             [(head, size, True, CYAN)])
        text(s, x + Inches(2.6), yy, w - Inches(2.6), Inches(0.5),
             [(body, size, False, DIM)])


prs = new_deck()

# ── S1 Title ────────────────────────────────────────────────────────────────
s = slide(prs)
text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2),
     [("NETRA", 72, True, INK),
      ("Networked Evidence, Tracking & Risk Analytics", 20, False, CYAN)])
text(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(1),
     [("From lakhs of FIRs to Monday-morning decisions.", 24, False, DIM)])
text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1),
     [("Challenge 2  ·  KSP Datathon 2026  ·  [TEAM NAME]", 14, False, MUTE)])
footer(s)
note(s, "Place the NETRA logo top-left or centred above the wordmark. "
        "Keep this slide sparse — no screenshot.")

# ── S2 Problem ──────────────────────────────────────────────────────────────
s = slide(prs)
heading(s, "Lakhs of FIRs. Almost no intelligence.")
accent_bar(s)
bullets(s, Inches(0.7), Inches(2.1), Inches(12), [
    ("Siloed", "1,100+ police stations file into one database nobody can see across."),
    ("Invisible serial crime", "One offender hitting 3 districts looks like 3 unrelated small cases, to 3 officers who never speak."),
    ("Always backwards", "Officers see last month's counts, never next week's risk."),
], size=15, gap=0.95)
text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(1),
     [("[ NCRB statistic — Karnataka case volume / pendency — Aditya to supply ]", 13, False, AMBER)])
footer(s)
note(s, "Do NOT invent the NCRB numbers. Leave the amber box until Aditya sends them.\n"
        "Optional diagram (mermaid in PPT_HANDBOOK S2): 3 stations -> 'no shared view' -> "
        "'3 small unsolved cases, 1 offender nobody sees'.")

# ── S3 Solution ─────────────────────────────────────────────────────────────
s = slide(prs)
heading(s, "One platform. Five capabilities.",
        "Built field-for-field on KSP's official FIR schema.")
accent_bar(s)
bullets(s, Inches(0.7), Inches(2.2), Inches(12), [
    ("SEE  ·  Map", "Live Karnataka map — where crime clusters, state → district → station"),
    ("TRACK  ·  Trends", "Trends over time + automatic anomaly alerts"),
    ("CONNECT  ·  Linkage", "Links the same offender across districts by method — English & ಕನ್ನಡ"),
    ("JUDGE  ·  Outcomes", "Case outcomes: chargesheeted / false / undetected"),
    ("ACT  ·  Forecast", "7-day forecast + patrol windows + daily briefing PDF"),
], size=14, gap=0.72)
text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
     [("CONNECT is our differentiator — make it visually heaviest.", 12, True, CYAN)])
footer(s)
note(s, "Replace these rows with the radial hub diagram (prompt in PPT_HANDBOOK S3): "
        "NETRA eye at centre, 5 nodes around it, CONNECT emphasised.")

# ── S4 Linkage (star) ───────────────────────────────────────────────────────
s = slide(prs)
heading(s, "One offender. Three districts. Found in under a second.",
        "AI reads the narrative, not keywords — in English and ಕನ್ನಡ.")
accent_bar(s)
picture(s, "S4-linkage-star.png", Inches(0.7), Inches(1.95), Inches(8.4))
text(s, Inches(9.4), Inches(2.0), Inches(3.4), Inches(3),
     [("THE PROOF", 11, True, CYAN),
      ("Blind test: NETRA recovered 5 of 5 hidden serial-crime patterns at 0.96 precision, "
       "mixed among 1,500 decoy FIRs.", 14, True, INK),
      ("The patterns were authored by a teammate and never shown to the person who built the AI.",
       11, False, DIM)])
text(s, Inches(9.4), Inches(4.6), Inches(3.4), Inches(2.2),
     [("Crime DNA — the shared modus-operandi fingerprint", 11, False, DIM),
      ("Predicted base zone — Rossmo profiling. A place, not a person.", 11, False, DIM),
      ("Next strike window — from the series' own cadence", 11, False, DIM)], spacing=1.4)
text(s, Inches(0.7), Inches(6.55), Inches(8.4), Inches(0.5),
     [("8 burglaries. 3 districts. 3 officers who never spoke. One shared method — surfaced automatically.",
       12, False, DIM)])
footer(s)
note(s, "THE star slide — give it the most airtime.\n"
        "Cohesion % = measured cosine similarity of shared MO, not a made-up confidence.\n"
        "Never say 'identifies suspects'. It predicts places and patterns.")

# ── S5 Forecast ─────────────────────────────────────────────────────────────
s = slide(prs)
heading(s, "Next week's hotspots — validated on real data.",
        "Tested on 503,468 real crimes, not just our own synthetic set.")
accent_bar(s)
picture(s, "S5-analytics-trends.png", Inches(0.7), Inches(1.95), Inches(7.6))
text(s, Inches(8.6), Inches(2.0), Inches(4.2), Inches(4.5),
     [("VALIDATION", 11, True, CYAN),
      ("503,468 real City-of-Chicago crimes (2022–23, 77 areas)", 12, False, INK),
      ("LSTM: top-10 hotspot hit-rate 0.867 · R² 0.950", 13, True, OK),
      ("Beats gradient boosting (0.858 / 0.948) — both on real data", 11, False, DIM),
      ("Honest: a moving-average baseline scores 0.85 on the same metric and better raw error "
       "(MAE 8.22 vs 8.71). We win on hotspot ranking — what a patrol planner uses — and do not "
       "claim to beat it on count error.", 10, False, AMBER)], spacing=1.3)
text(s, Inches(8.6), Inches(5.5), Inches(4.2), Inches(1.5),
     [("PATROL OPTIMIZER", 11, True, CYAN),
      ("Submodular allocation · provable (1−1/e) ≈ 63% optimality. Deterrence grounded in the "
       "Minneapolis Hot-Spots RCT and the Koper Curve. Top 10 of 77 areas hold 33.6% of crime.",
       10, False, DIM)], spacing=1.25)
footer(s)
note(s, "KEEP the amber honesty line. Publishing where the baseline wins is what makes the "
        "rest believable to a technical judge.\n"
        "Optional chart (PPT_HANDBOOK S5): grouped horizontal bar, LSTM vs GBM on hit@10 / R² / MAE. "
        "Y-axis must start at 0.")

# ── S6 Product (2x2) ────────────────────────────────────────────────────────
s = slide(prs)
heading(s, "Not slides — a working system.", "Deployed and live on Zoho Catalyst.")
accent_bar(s)
# A single row of four: 16:9 screenshots stacked two-high overflow the slide height.
grid = [("S6-command-map.png", "Command Map", "5,000 incidents, hotspot density, drill-down"),
        ("S6e-offender-network.png", "Offender network", "rings, kingpins, predicted ties"),
        ("S6b-alerts.png", "Anomaly alerts", "each with a why, and a lifecycle to work"),
        ("S6f-case-outcomes.png", "Case outcomes", "chargesheeted / false / undetected")]
IMG_W = 2.95
GAP = 0.19
span = 4 * IMG_W + 3 * GAP
x0 = (13.333 - span) / 2
for i, (img, cap, sub) in enumerate(grid):
    x = Inches(x0 + i * (IMG_W + GAP))
    picture(s, img, x, Inches(2.55), Inches(IMG_W))
    text(s, x, Inches(4.45), Inches(IMG_W), Inches(0.8),
         [(cap, 11, True, CYAN), (sub, 9, False, DIM)], spacing=1.2)
text(s, Inches(0.7), Inches(5.75), Inches(12), Inches(0.5),
     [("Every number on every screen is computed from the data — nothing here is mocked-up art.",
       12, False, MUTE)], align=PP_ALIGN.CENTER)
note(s, "Spares if a tile doesn't fit: S6c-case-search.png (semantic MO search), "
        "S6d-documents.png (document centre).")

# ── S7 Briefing ─────────────────────────────────────────────────────────────
s = slide(prs)
heading(s, "Every SP's Monday morning, written for them.",
        "English and ಕನ್ನಡ · one page · exports to PDF.")
accent_bar(s)
picture(s, "S7-briefing.png", Inches(0.7), Inches(1.95), Inches(7.8))
text(s, Inches(8.8), Inches(2.1), Inches(4.0), Inches(4),
     [("Auto-written per district from that district's own numbers — never a template with blanks filled in.",
       12, False, DIM),
      ("Kannada is first-class, not an afterthought — the same briefing, natively.", 12, False, DIM),
      ("Exports as a police-letterhead PDF, plus a fax-ready monochrome sheet for stations still on fax.",
       12, False, DIM)], spacing=1.5)
text(s, Inches(8.8), Inches(5.8), Inches(4.0), Inches(0.8),
     [("Decision support — verify against ground reports before acting (human-in-the-loop).",
       10, False, MUTE)])
footer(s)

# ── S8 Schema ───────────────────────────────────────────────────────────────
s = slide(prs)
heading(s, "This plugs into real data on day one.",
        "Field-for-field faithful to KSP's official FIR ER model.")
accent_bar(s)
text(s, Inches(0.7), Inches(2.2), Inches(5.9), Inches(3.5),
     [("[ KSP official ER-diagram crop — Aditya to supply ]", 13, False, AMBER)])
text(s, Inches(6.9), Inches(2.2), Inches(5.9), Inches(3.5),
     [("Our tables", 13, True, CYAN),
      ("Mirror KSP's ER model 1:1 — same entities, same relationships. See docs/SYSTEM_DESIGN.md §3.",
       12, False, DIM)], spacing=1.3)
text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1),
     [("We did not invent a convenient schema. We mirrored theirs — so this is an integration, not a rebuild.",
       14, True, INK),
      ("All demo records are synthetic, generated TO that schema. KSP's real FIR data is confidential; "
       "only the schema is public.", 11, False, AMBER)], spacing=1.3)

# ── S9 Architecture ─────────────────────────────────────────────────────────
s = slide(prs)
heading(s, "Runs entirely on Zoho Catalyst.",
        "Sovereign by design — FIR text never leaves the device.")
accent_bar(s)
text(s, Inches(0.7), Inches(2.1), Inches(6.4), Inches(4),
     [("[ Architecture diagram — mermaid source in PPT_HANDBOOK S9 ]", 13, False, AMBER),
      ("Browser: React SPA + sentence-transformer running IN-BROWSER", 11, False, DIM),
      ("Catalyst: Web Client Hosting · Advanced I/O Function (netra_api) · Data Store · QuickML",
       11, False, DIM),
      ("Offline: Python pipeline — clusters, forecasts, risk, network", 11, False, DIM)], spacing=1.4)
bullets(s, Inches(7.4), Inches(2.2), Inches(5.4), [
    ("Web Client Hosting", "Serves the React app"),
    ("Advanced I/O Function", "netra_api — Express API, RBAC, LLM proxy"),
    ("Data Store", "Relational — KSP FIR schema + analytics"),
    ("QuickML", "GLM-4.7 text + Qwen VLM for scanned docs"),
], size=11, gap=0.62)
text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.8),
     [("Semantic matching runs in the browser, so FIR narratives never reach any external inference API. "
       "Model weights still load from the open-weights host on first use — documented in docs/SOVEREIGN-MODEL.md.",
       10, False, MUTE)])
note(s, "Use the diagram in PPT_HANDBOOK S9 — NOT the older AppSail diagram in SYSTEM_DESIGN.md, "
        "which describes a target design we did not ship. Do not overclaim 'fully sovereign'.")

# ── S10 Responsible AI ──────────────────────────────────────────────────────
s = slide(prs)
heading(s, "The constraints came first.", "Not a compliance slide — these are enforced in code.")
accent_bar(s)
bullets(s, Inches(0.7), Inches(2.2), Inches(12), [
    ("No caste or religion", "Those fields exist in the schema for fidelity, but are excluded from every model and never rendered by any screen or API."),
    ("Places, never people", "Forecast features are place / time / crime-type only. The geographic profile predicts a zone, not a person."),
    ("Every output explains itself", "Each alert carries a why-flagged statistic; each AI answer carries citations and a reasoning trace."),
    ("A human always decides", "Every screen and exported document is labelled decision support, human-in-the-loop."),
], size=13, gap=0.95)
text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.6),
     [("An officer can always ask \"why did it say that?\" — and get an answer.", 15, True, CYAN)])
footer(s)

# ── S11 Honest engineering ──────────────────────────────────────────────────
s = slide(prs)
heading(s, "We tested it like we expect to be doubted.",
        "Showing the limits is what makes the numbers believable.")
accent_bar(s)
bullets(s, Inches(0.7), Inches(2.3), Inches(12), [
    ("Blind test", "5 patterns authored by a teammate, hidden among 1,500 decoys, answer key withheld until after the run. Result: 5/5 recovered, 0.96 precision."),
    ("Benchmarked honestly", "Measured against a strong baseline on real data — and we publish where the baseline wins. A team that only reports wins is hiding something."),
    ("Limits in the repo", "docs/SECURITY-NOTES.md and docs/SOVEREIGN-MODEL.md record what is accepted, unfinished, or constrained by the platform."),
], size=13, gap=1.15)
footer(s)
note(s, "Optional image: the data-quality panel (e.g. '94.1% geocoded') — showing coverage gaps "
        "rather than hiding them. Ask Aditya to capture it.")

# ── S12 Impact & roadmap ────────────────────────────────────────────────────
s = slide(prs)
heading(s, "From prototype to statewide.",
        "Built by 3 students, ₹0 infrastructure cost, in ~2 weeks.")
accent_bar(s)
bullets(s, Inches(0.7), Inches(2.1), Inches(12), [
    ("One arrest, many closures", "The Linkage screen names how many unsolved FIRs a single arrest could clear."),
    ("Patrol where it matters", "Top 10 of 77 areas held 33.6% of crime on real data — concentration is measurable, not a hunch."),
    ("Minutes, not weeks", "Cross-district serial crime goes from 'nobody noticed' to under a second."),
], size=13, gap=0.85)
for i, (when, what) in enumerate([
        ("NOW — Prototype", "Synthetic data on KSP's schema, deployed on Catalyst"),
        ("NEXT — Pilot", "One district, real FIR data behind KSP's firewall, officer feedback loop"),
        ("THEN — Statewide", "All districts, live ingestion, a briefing to every SP each morning")]):
    x = Inches(0.7 + i * 4.15)
    text(s, x, Inches(4.9), Inches(3.9), Inches(1.2),
         [(when, 13, True, CYAN), (what, 11, False, DIM)], spacing=1.25)
text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
     [("NETRA doesn't replace an officer's judgement. It makes sure nothing reaches them too late.",
       15, True, INK)])
footer(s)

prs.save(OUT)
print(f"✓ {OUT.name} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
